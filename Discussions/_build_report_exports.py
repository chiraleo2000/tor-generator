"""Build DOCX, PPTX, and PDF copies of report 19."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRgb
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn as ppt_qn
from pptx.util import Inches as PptInches, Pt as PptPt

ROOT = Path(__file__).resolve().parent
EVIDENCE = ROOT / "test-evidence"
OUT_DOCX = ROOT / "19-APPLICATION_OPERATING_REPORT.docx"
OUT_PPTX = ROOT / "19-APPLICATION_OPERATING_REPORT.pptx"
OUT_PDF = ROOT / "19-APPLICATION_OPERATING_REPORT.pdf"
FONT_DIR = Path.home() / "AppData/Local/Microsoft/Windows/Fonts"
FONT_REG = FONT_DIR / "THSarabunNew.ttf"
FONT_BOLD = FONT_DIR / "THSarabunNew Bold.ttf"
FONT_NAME = "TH Sarabun New"

NAVY = RGBColor(0x1D, 0x40, 0x8C)
NAVY_DARK = RGBColor(0x0F, 0x23, 0x3E)
ORANGE = RGBColor(0xF2, 0xA3, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5B, 0x67, 0x75)
HEADER_FILL = "1D408C"

P_NAVY = PptRgb(0x1D, 0x40, 0x8C)
P_NAVY_DARK = PptRgb(0x0F, 0x23, 0x3E)
P_ORANGE = PptRgb(0xF2, 0xA3, 0x0A)
P_WHITE = PptRgb(0xFF, 0xFF, 0xFF)
P_BODY = PptRgb(0x1F, 0x2A, 0x37)
P_MUTED = PptRgb(0x5B, 0x67, 0x75)
P_ROW = PptRgb(0xF4, 0xF7, 0xFB)


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_BOLD if bold and FONT_BOLD.exists() else FONT_REG
    return ImageFont.truetype(str(path), size)


def draw_architecture() -> Path:
    w, h = 1400, 620
    img = Image.new("RGB", (w, h), (255, 255, 255))
    d = ImageDraw.Draw(img)
    title = _font(28, True)
    body = _font(18)
    small = _font(16)
    d.text((40, 18), "สถาปัตยกรรมระบบ TOR Generator", font=title, fill=(15, 35, 62))
    boxes = [
        (40, 80, 220, 170, "ผู้ใช้\nเบราว์เซอร์", (242, 163, 10)),
        (320, 80, 560, 170, "Next.js :3000\nrewrite /api/v1", (29, 64, 140)),
        (660, 80, 920, 170, "FastAPI :4000\n/api/v1", (29, 64, 140)),
    ]
    for x1, y1, x2, y2, text, color in boxes:
        d.rounded_rectangle([x1, y1, x2, y2], radius=12, fill=color)
        d.text((x1 + 16, y1 + 22), text, font=body, fill=(255, 255, 255))
    d.polygon([(232, 125), (308, 115), (308, 135)], fill=(91, 103, 117))
    d.polygon([(572, 125), (648, 115), (648, 135)], fill=(91, 103, 117))
    stores = [
        (40, 250, "PostgreSQL\n+ pgvector"),
        (240, 250, "MongoDB\nGridFS"),
        (440, 250, "Neo4j\nGraphRAG"),
        (640, 250, "Redis"),
        (800, 250, "MinIO"),
        (960, 250, "LLM /\nEmbeddings"),
        (1160, 250, "Rule\nEngine"),
    ]
    for x, y, text in stores:
        d.rounded_rectangle([x, y, x + 170, y + 110], radius=10, fill=(15, 35, 62))
        d.text((x + 14, y + 28), text, font=small, fill=(255, 255, 255))
    d.text(
        (40, 400),
        "เบราว์เซอร์เรียก /api/v1 โดเมนเดียวกับ UI  →  Next.js rewrite ไป backend:4000 (timeout 5 นาที)",
        font=body,
        fill=(31, 42, 55),
    )
    d.text(
        (40, 450),
        "ค่าเริ่มต้น: Gemma ใน LM Studio + EmbeddingGemma 768-d  ·  แชทและ embeddings เลือกอิสระ",
        font=body,
        fill=(31, 42, 55),
    )
    d.text(
        (40, 500),
        "JWT คุกกี้ HttpOnly  ·  บทบาท officer / reviewer / admin  ·  คลังกฎหมายจาก PDF ผ่าน seed_raw_docs",
        font=body,
        fill=(31, 42, 55),
    )
    d.text(
        (40, 560),
        "v0.2.3  ·  21 สิงหาคม 2026",
        font=small,
        fill=(91, 103, 117),
    )
    path = EVIDENCE / "19-diagram-architecture.png"
    img.save(path, "PNG")
    return path


def draw_phases() -> Path:
    w, h = 1400, 320
    img = Image.new("RGB", (w, h), (255, 255, 255))
    d = ImageDraw.Draw(img)
    title = _font(26, True)
    body = _font(18, True)
    small = _font(16)
    d.text((40, 16), "พื้นที่ทำงาน 5 Phase (เส้นทางหลักบน UI)", font=title, fill=(15, 35, 62))
    labels = [
        ("0", "อัปโหลด", "ชุดเอกสาร + แชท"),
        ("1", "ช่องว่าง", "coverage + พร้อมร่าง"),
        ("2", "ร่าง", "s1–s13 + HITL"),
        ("3", "ตรวจ", "Rule Engine ≥ 70"),
        ("4", "ส่งออก", "Word / PDF"),
    ]
    x = 40
    for i, (num, name, hint) in enumerate(labels):
        d.rounded_rectangle([x, 80, x + 220, 250], radius=12, fill=(29, 64, 140))
        d.ellipse([x + 16, 96, x + 56, 136], fill=(242, 163, 10))
        d.text((x + 28, 100), num, font=body, fill=(15, 35, 62))
        d.text((x + 70, 102), name, font=body, fill=(255, 255, 255))
        d.text((x + 16, 160), hint, font=small, fill=(230, 235, 245))
        if i < len(labels) - 1:
            d.polygon([(x + 228, 155), (x + 252, 145), (x + 252, 165)], fill=(242, 163, 10))
        x += 268
    path = EVIDENCE / "19-diagram-phases.png"
    img.save(path, "PNG")
    return path


def _set_run_font(run, size_pt: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    for key in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        rfonts.set(qn(key), FONT_NAME)


def _shade(cell, fill: str) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    shd.set(qn("w:val"), "clear")
    tc_pr.append(shd)


def _set_cell_text(cell, text: str, *, bold: bool = False, color: RGBColor | None = None, size: int = 12) -> None:
    cell.text = ""
    para = cell.paragraphs[0]
    para.paragraph_format.space_after = Pt(0)
    run = para.add_run(text)
    _set_run_font(run, size, bold=bold, color=color or BODY)


def add_table(doc: Document, headers: list[str], rows: list[list[str]]) -> None:
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        _shade(cell, HEADER_FILL)
        _set_cell_text(cell, header, bold=True, color=WHITE, size=12)
    for r_idx, row in enumerate(rows):
        fill = "F4F7FB" if r_idx % 2 else "FFFFFF"
        for c_idx, value in enumerate(row):
            cell = table.rows[r_idx + 1].cells[c_idx]
            _shade(cell, fill)
            _set_cell_text(cell, value, size=12)
    doc.add_paragraph()


def add_heading(doc: Document, text: str, level: int) -> None:
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        _set_run_font(run, 18 if level == 1 else 16 if level == 2 else 14, bold=True, color=NAVY)


def add_body(doc: Document, text: str, *, italic: bool = False) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run(text)
    _set_run_font(run, 14, color=BODY)
    run.italic = italic


def add_image(doc: Document, path: Path, caption: str, width_cm: float = 16.0) -> None:
    if not path.exists():
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(path), width=Cm(width_cm))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cap.add_run(caption)
    _set_run_font(r, 12, color=MUTED)
    r.italic = True


def add_page_numbers(doc: Document) -> None:
    footer = doc.sections[0].footer
    para = footer.paragraphs[0]
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = para.add_run("หน้า ")
    _set_run_font(run, 12, color=MUTED)
    fld = para.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    fld._r.append(begin)
    instr_run = para.add_run()
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    instr_run._r.append(instr)
    end_run = para.add_run()
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    end_run._r.append(end)


def build_docx(arch: Path, phases: Path) -> None:
    doc = Document()
    section = doc.sections[0]
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    for attr in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(section, attr, Cm(2.2))
    style = doc.styles["Normal"]
    style.font.name = FONT_NAME
    style.font.size = Pt(14)
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    for key in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        rfonts.set(qn(key), FONT_NAME)

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    hr = header.add_run("TOR Generator  v0.2.3  ·  รายงานการทำงานของแอปพลิเคชัน")
    _set_run_font(hr, 11, color=MUTED)
    add_page_numbers(doc)

    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("รายงานการทำงานของแอปพลิเคชัน")
    _set_run_font(r, 22, bold=True, color=NAVY)
    t2 = doc.add_paragraph()
    t2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = t2.add_run("TOR Generator")
    _set_run_font(r2, 28, bold=True, color=NAVY_DARK)
    t3 = doc.add_paragraph()
    t3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = t3.add_run("ระบบจัดทำและตรวจสอบ TOR ตาม พ.ร.บ. การจัดซื้อจัดจ้างฯ พ.ศ. 2560")
    _set_run_font(r3, 16, color=BODY)
    t4 = doc.add_paragraph()
    t4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r4 = t4.add_run("เวอร์ชัน 0.2.3  ·  วันที่ 21 สิงหาคม 2026 (พ.ศ. 2569)")
    _set_run_font(r4, 14, color=MUTED)
    add_body(
        doc,
        "แหล่งความจริง: โค้ดใน app/frontend และ app/backend รวมผล unit tests ที่รันจริงวันเดียวกัน "
        "(Vitest 167 ผ่าน, pytest 1500 + live_llm 14, Playwright headed 20 · พิมพ์ช้า + LM Studio จริง) เอกสารนี้เป็นฉบับส่งออกของ Discussions/19-APPLICATION_OPERATING_REPORT.md",
    )

    add_heading(doc, "1. สรุปหนึ่งหน้า", 1)
    add_body(
        doc,
        "แอปนี้ช่วยเจ้าหน้าที่พัสดุร่างและตรวจ TOR โดยบังคับโครง 13 ส่วน (s1–s13) รับชุดเอกสาร "
        "จัดเข้าช่อง ตรวจด้วย Rule Engine ยืนยันหมวดเสี่ยงด้วยคน แล้วส่งออก Word/PDF รูปแบบราชการ",
    )
    add_table(
        doc,
        ["ข้อ", "ค่า"],
        [
            ["ปัญหาที่แก้", "เขียน TOR ด้วยมือช้า ตกหัวข้อกฎหมาย จัดรูปแบบราชการยาก"],
            ["UI", "Next.js 14 App Router · พอร์ต 3000 · ภาษาไทยทั้งหน้าจอ"],
            ["API", "FastAPI · พอร์ต 4000 · คำนำหน้า /api/v1"],
            ["พื้นที่ทำงานหลัก", "5 Phase (0–4) ไม่ใช่วิซาร์ด 8 ขั้น"],
            ["โมเดลเอกสาร", "13 ส่วนกฎหมาย + s4.1–s4.14 = 27 ช่อง"],
            ["AI ค่าเริ่มต้น", "Gemma ใน LM Studio + EmbeddingGemma 768 มิติ"],
            ["บทบาท", "officer / reviewer / admin"],
            ["เทสต์รอบนี้", "Vitest 167 · pytest 1500 + live_llm 14 · Playwright headed 20 · cov 86%/87.09%"],
        ],
    )
    add_image(doc, arch, "ภาพที่ 1  สถาปัตยกรรมรวม — ผู้ใช้ → Next.js → FastAPI → คลังข้อมูล / LLM / Rule Engine")

    add_heading(doc, "2. ผู้ใช้ สิทธิ์ และหน้าจอหลัก", 1)
    add_table(
        doc,
        ["บทบาท", "ทำได้", "ทำไม่ได้"],
        [
            ["เจ้าหน้าที่ (officer)", "สร้าง/แก้ไขโครงการของตนเอง · คลังส่วนตัว · ถาม-ตอบ · ส่งขออนุมัติ", "เห็นโครงการคนอื่น · อนุมัติ · ตั้งค่า AI ระบบ"],
            ["ผู้ตรวจสอบ (reviewer)", "เห็นทุกโครงการ · อนุมัติ / ส่งกลับ", "จัดการแม่แบบ ผู้ใช้ การตั้งค่า AI"],
            ["ผู้ดูแล (admin)", "ทุกอย่างของ reviewer + แม่แบบ คลังกลาง ผู้ใช้ ตั้งค่า AI", "—"],
        ],
    )
    add_body(doc, "JWT อยู่ในคุกกี้ HttpOnly tor_access_token (SameSite=Lax) และรองรับ Authorization: Bearer สำหรับเทสต์และไคลเอนต์ API")
    add_table(
        doc,
        ["เมนู", "เส้นทาง", "ใช้ทำอะไร"],
        [
            ["แดชบอร์ด", "/projects", "รายการโครงการ สร้างใหม่ อนุมัติ/ส่งกลับ"],
            ["ฐานความรู้", "/knowledge-base", "คลังกลาง + เอกสารของฉัน"],
            ["ร่าง TOR", "/projects/{id}/draft", "พื้นที่ทำงาน 5 Phase"],
            ["ตรวจสอบ TOR", "/review", "ตรวจไฟล์ภายนอก + เทียบ Jaccard"],
            ["ถาม-ตอบ", "/chat", "ห้องแชทคลังความรู้ (SSE)"],
            ["คู่มือ", "/help", "FAQ และการใช้งาน"],
            ["แอดมิน", "/admin/...", "แม่แบบ ผู้ใช้ คลังกลาง ตั้งค่า AI"],
        ],
    )
    add_image(doc, EVIDENCE / "00-login.png", "ภาพที่ 2  หน้าเข้าสู่ระบบ")
    add_image(doc, EVIDENCE / "02-dashboard.png", "ภาพที่ 3  แดชบอร์ดหลังล็อกอิน")
    add_heading(doc, "2.1 สามเครื่องมือหลักของเจ้าหน้าที่", 2)
    add_table(
        doc,
        ["เครื่องมือ", "หน้า", "ทำอะไร", "ผลที่ควรเห็น"],
        [
            ["ร่าง TOR", "/projects/{id}/draft", "5 Phase: อัปโหลด → เติมช่อง → ร่าง 13 หมวด (+ AI) → HITL", "มีเนื้อหา s1–s13 พร้อมส่งตรวจ"],
            ["ตรวจสอบ TOR", "Phase 4 + /review", "Rule Engine ≥ 70 + ReviewAgent · ตรวจไฟล์ภายนอก", "คะแนน / findings / suggestions"],
            ["ถาม-ตอบ", "/chat", "ห้องคลังความรู้ SSE + citations จาก RAG", "คำตอบยาวพร้อมชิปอ้างอิง"],
        ],
    )

    add_heading(doc, "3. Frontend — Next.js 14", 1)
    add_body(
        doc,
        "โฟลเดอร์ app/frontend/src แยกเป็นหน้า App Router, คอมโพเนนต์, lib และ Zustand stores "
        "เบราว์เซอร์เรียก /api/v1 โดเมนเดียวกับ UI แล้ว Next.js rewrite ไป backend:4000 timeout 5 นาที "
        "apiClient ส่งคุกกี้เสมอ ถ้า HTTP 401 จะล้างเซสชันแล้วไป /login",
    )
    add_table(
        doc,
        ["กลุ่ม", "หน้า", "การ์ด"],
        [
            ["สาธารณะ", "/", "ส่งแขกไป /login"],
            ["ยืนยันตัว", "/login /register", "เลย์เอาต์การ์ดกรมทัพเรือ"],
            ["งานหลัก", "/projects /draft /chat /knowledge-base /review /help", "AuthGuard + ไซด์บาร์ 255px"],
            ["แอดมิน", "/admin/templates users knowledge-base ai-settings", "ถ้าไม่ใช่แอดมินจะถูกเด้ง"],
        ],
    )
    add_table(
        doc,
        ["Store", "หน้าที่"],
        [
            ["auth-store", "ผู้ใช้ + โทเค็นในหน่วยความจำ · restoreSession เรียก GET /auth/me · ไม่เก็บ JWT ใน localStorage"],
            ["project-store", "รายการโครงการ โครงการที่เปิดอยู่ สร้าง/อัปเดต/ส่งตรวจ"],
            ["wizard-store", "เหลือสำหรับวิซาร์ดเก่า — ตัวร่างจริงคือ DraftWorkspace"],
            ["ui-store", "ธีม ไซด์บาร์ โทสต์"],
        ],
    )
    add_table(
        doc,
        ["คอมโพเนนต์", "บทบาท"],
        [
            ["DraftWorkspace", "พื้นที่ทำงาน Phase 0–4 ทั้งก้อน"],
            ["IntakeChatPanel", "แชทร่าง + อัปโหลดชุดเอกสาร (Phase 0–2)"],
            ["ChatShell + MiniRoomList", "โครงห้องแชทแบบย่อ (ถาม-ตอบคลัง)"],
            ["PhaseFlow", "แถบ 5 Phase + ล็อกตามเกต"],
            ["NewProjectDialog", "สร้างโครงการ (ชื่อ หน่วยงาน วงเงิน ASCII ประเภท แม่แบบ)"],
            ["ProjectRowActions", "เปิดร่าง / อนุมัติ / ส่งกลับ"],
        ],
    )
    add_body(doc, "วิซาร์ด 8 ขั้นยังอยู่ในโค้ดเพื่อความเข้ากันได้ — หน้าจอที่ใช้งานคือ 5 Phase")

    add_heading(doc, "4. Backend — FastAPI + LangGraph", 1)
    add_table(
        doc,
        ["ชั้น", "โฟลเดอร์", "ทำอะไร"],
        [
            ["HTTP", "app/api/v1/endpoints/", "REST + SSE"],
            ["Domain", "app/domain/", "s1–s13, ช่อง intake, magic bytes"],
            ["Orchestrator", "app/orchestrator/", "กราฟร่างรายหมวด + กราฟเอเจนต์ทั้งฉบับ"],
            ["Agents", "app/orchestrator/agents/", "เอเจนต์ 13 หมวด + ReviewAgent"],
            ["RAG", "app/rag/", "สกัด หั่น ingest ค้น กราฟ ACL"],
            ["Rules", "app/rule_engine/", "คะแนนคุณภาพแบบกำหนดได้ซ้ำ"],
            ["Providers", "app/providers/", "LLM / embeddings / vector store"],
            ["Export", "app/export/", "DOCX PDF รูปแบบไทย MinIO"],
            ["Services", "app/services/", "intake, coverage, gap, auth, audit, agent"],
            ["Models", "app/models/", "SQLAlchemy"],
        ],
    )
    add_heading(doc, "4.1 โมเดล TOR 13 ส่วน", 2)
    add_table(
        doc,
        ["คีย์", "ชื่อไทย", "เอเจนต์", "HITL"],
        [
            ["s1", "ความเป็นมา", "Background", "ไม่"],
            ["s2", "วัตถุประสงค์", "Objectives", "ไม่"],
            ["s3", "คุณสมบัติผู้เสนอราคา", "Qualifications", "ใช่"],
            ["s4", "ขอบเขตของงาน (+ s4.1–s4.14)", "Scope", "ไม่"],
            ["s5", "ระยะเวลาดำเนินการ", "Timeline", "ไม่"],
            ["s6", "วงเงินงบประมาณ", "Budget", "ใช่"],
            ["s7", "สถานที่ดำเนินการ", "Location", "ไม่"],
            ["s8", "งวดงานและการจ่ายเงิน", "Payment", "ใช่"],
            ["s9", "การรับประกัน", "Warranty", "ไม่"],
            ["s10", "อัตราค่าปรับ", "Penalties", "ใช่"],
            ["s11", "เกณฑ์พิจารณาข้อเสนอ", "Evaluation", "ไม่"],
            ["s12", "เอกสารที่ต้องยื่น", "Documents", "ไม่"],
            ["s13", "เงื่อนไขอื่น ๆ", "Conditions", "ใช่"],
        ],
    )
    add_body(doc, "หมวด HITL (s3 s6 s8 s10 s13) ต้องกดยืนยันคนก่อนปุ่มส่งขออนุมัติจะเปิด  ขอบเขตงานย่อยที่บังคับขั้นต่ำ: s4.1 สรุปขอบเขต และ s4.8 ผลงานส่งมอบ")

    add_heading(doc, "5. Workflows", 1)
    add_heading(doc, "5.1 พื้นที่ทำงาน 5 Phase", 2)
    add_image(doc, phases, "ภาพที่ 4  เส้นทางหลัก Phase 0 → 4")
    add_table(
        doc,
        ["Phase", "ผู้ใช้ทำอะไร", "API หลัก", "เกต"],
        [
            ["0", "สร้างโครงการ วางข้อความหรืออัปโหลด แล้วกดเริ่มวิเคราะห์", "intake/upload text analyze", "มีเนื้อหา + กดวิเคราะห์ → ปลด Phase 1"],
            ["1", "ตารางความครบ ดึงกฎหมายครั้งเดียว นับ 10 วินาทีหรือกดไปเลย", "coverage fill-references", "วิเคราะห์แล้ว → ปลด Phase 2 (ไม่มีไดอะล็อก)"],
            ["2", "ตารางคู่แชทถามช่องที่ขาด ดึงอ้างอิงจากชิป แล้วกดพร้อมร่าง", "chat fill-reference confirm-ready", "ready_to_compose → ปลด Phase 3"],
            ["3", "ร่าง 13 หมวดอัตโนมัติ ยืนยัน HITL แล้วไปทบทวน", "sections draft-section draft-chat confirm-phase4", "ครบ 13 หมวดร่างแล้ว → ปลด Phase 4"],
            ["4", "แชทรีวิว Rule Engine อัตโนมัติ ส่งออก DOCX/PDF", "review submit export", "reviewer/admin อนุมัติหรือส่งกลับ"],
        ],
    )
    add_image(doc, EVIDENCE / "03-phase-0-upload.png", "ภาพที่ 5  Phase 0 อัปโหลดชุดเอกสาร")
    add_image(doc, EVIDENCE / "04b-phase-1-coverage.png", "ภาพที่ 6  Phase 1 ตารางความครบ")
    add_image(doc, EVIDENCE / "05-phase-2-draft.png", "ภาพที่ 7  Phase 3 ร่าง 13 หมวด")
    add_image(doc, EVIDENCE / "05b-hitl-confirm.png", "ภาพที่ 8  ยืนยัน HITL ก่อนส่งตรวจ")
    add_image(doc, EVIDENCE / "e2e-phase-4-review-chat.png", "ภาพที่ 9  Phase 4 แชทรีวิวและ Rule Engine")
    add_image(doc, EVIDENCE / "07-phase-4-publish.png", "ภาพที่ 10  Phase 4 ส่งออก Word/PDF")

    add_heading(doc, "5.2 กราฟร่างรายหมวด (LangGraph)", 2)
    add_body(
        doc,
        "ปุ่มร่างด้วย AI ใน Phase 3 เดิน validate_input → retrieve_context (RAG) → llm_draft (เอเจนต์ s1–s13) "
        "→ rule_guardrail (คะแนน ≥ 70 ผ่าน, ไม่ผ่าน retry สูงสุด 3 ครั้ง) → human_review → finalize  "
        "Timeout ต่อหมวดค่าเริ่มต้น 180 วินาทีสำหรับ LM Studio (สูงสุด 300)",
    )
    add_heading(doc, "5.3 เส้นทางเอเจนต์ทั้งฉบับ", 2)
    add_body(
        doc,
        "/api/v1/agent เป็นเส้นขนานยังไม่มีหน้า Next.js: ingest → แผนที่ 27 ช่อง → ถามส่วนขาด (สูงสุด 20 รอบ) "
        "→ ยืนยัน → ร่างทั้งฉบับ → HITL → ส่งออก  สถานะเก็บใน agent_sessions.graph_state + Redis ไม่มี checkpointer บน Postgres",
    )
    add_heading(doc, "5.4 ถาม-ตอบ ตรวจไฟล์ภายนอก และการอนุมัติ", 2)
    add_body(
        doc,
        "หน้า /chat ใช้ห้อง kind=kb บน /api/v1/chat (SSE + อ้างอิง + แนบไฟล์คลังส่วนตัว) คนละประวัติกับแชทร่าง kind=draft_intake  "
        "หน้า /review สกัดแล้วเทียบ Jaccard โดยไม่ต้องสร้างโครงการ  เจ้าหน้าที่ส่งตรวจแล้ว reviewer/admin กดอนุมัติหรือส่งกลับบนแดชบอร์ด",
    )
    add_image(doc, EVIDENCE / "13-kb-chat.png", "ภาพที่ 11  ถาม-ตอบคลังความรู้")
    add_image(doc, EVIDENCE / "12-standalone-review.png", "ภาพที่ 12  ตรวจไฟล์ TOR ภายนอก")

    add_heading(doc, "6. Features", 1)
    add_table(
        doc,
        ["ฟีเจอร์", "สิ่งที่เกิดขึ้น", "หมายเหตุ"],
        [
            ["สร้างโครงการ", "ฟอร์มชื่อ หน่วยงาน วงเงิน ASCII ประเภท แม่แบบ", "เข้า /projects/{id}/draft"],
            ["Intake หลายไฟล์", "PDF/DOCX/PPTX/TXT ตรวจ magic bytes", "ต้นฉบับ GridFS"],
            ["แผนที่ช่องอัตโนมัติ", "NLP จัดเข้า s1–s13 / s4.1–s4.14", "ตารางความครบ Phase 1"],
            ["แชทร่าง", "SSE ถามส่วนขาด ดึงอ้างอิงกฎหมาย", "ไม่ปนกับ /chat"],
            ["ร่างด้วย AI รายหมวด", "เอเจนต์เฉพาะทาง + RAG + Rule Engine", "กรอกมือได้ถ้า LLM ไม่เปิด"],
            ["HITL 5 หมวด", "ปุ่มยืนยันก่อนส่งตรวจ", "s3 s6 s8 s10 s13"],
            ["Rule Engine", "คะแนน 0–100 ผ่านที่ 70", "กฎหมาย 40 / ครบ 30 / สอดคล้อง 20 / รูปแบบ 10"],
            ["ส่งออก DOCX/PDF", "TH Sarabun New · วันที่ พ.ศ.", "MinIO"],
            ["คลังส่วนตัว / คลังกลาง", "mine ตาม owner_id / แอดมินอัปโหลด", "ACL กันคนอื่นเห็น"],
            ["ตั้งค่า AI ทันที", "PUT /admin/ai-settings", "แชทกับ embeddings เลือกอิสระ ไม่รีสตาร์ท"],
        ],
    )
    add_image(doc, EVIDENCE / "11-knowledge-base.png", "ภาพที่ 13  ฐานความรู้ของเจ้าหน้าที่")
    add_image(doc, EVIDENCE / "09-admin-ai-lm-studio.png", "ภาพที่ 14  การตั้งค่า AI")

    add_heading(doc, "7. Functions", 1)
    add_heading(doc, "7.1 Frontend (src/lib)", 2)
    add_table(
        doc,
        ["ฟังก์ชัน", "ไฟล์", "ใช้ตอน"],
        [
            ["validatePassword", "password-rules.ts", "สมัครสมาชิก — ยาว ≥8 พิมพ์ใหญ่/เล็ก ตัวเลข ASCII อักขระพิเศษ"],
            ["canSelectPhase", "phase-gate.ts", "ล็อกแถบ Phase ไม่ให้ข้ามไปร่าง"],
            ["streamSsePost", "chat-sse.ts", "สตรีมคำตอบแชทคลังและแชทร่าง"],
            ["applyRequestAuth", "api-client.ts", "แนบโทเค็น / เด้ง login เมื่อ 401"],
            ["unwrapData", "api-unwrap.ts", "แปลงซอง { ok, data } เป็นอ็อบเจ็กต์ UI"],
            ["toReviewFinding", "review-findings.ts", "แผนที่ข้อค้นพบ Rule Engine"],
            ["jaccard / compareExtractJobs", "review-compare.ts", "เทียบ TOR ภายนอก"],
            ["llmOptionsForMode", "ai-settings.ts", "โหมดไม่บังคับสลับฝั่ง embeddings"],
        ],
    )
    add_heading(doc, "7.2 Backend", 2)
    add_table(
        doc,
        ["ฟังก์ชัน", "ที่อยู่", "ใช้ตอน"],
        [
            ["get_agent_for_section", "agents/registry.py", "จับ s1–s13 ไปเอเจนต์"],
            ["compile_agent_workflow_graph", "agent_graph.py", "กราฟเอเจนต์ทั้งฉบับ"],
            ["กราฟรายหมวด", "orchestrator/graph.py", "ปุ่มร่างด้วย AI"],
            ["RuleEngine.validate", "rule_engine/engine.py", "คะแนนถ่วงน้ำหนัก + ข้อค้นพบ"],
            ["ProviderFactory", "providers/factory.py", "เลือก LLM / embeddings / vector อิสระ"],
            ["require_role / require_project_access", "rbac.py", "กัน API ตามบทบาทและเจ้าของ"],
            ["compute_ready", "services/coverage.py", "ช่องครบพอจะยืนยันพร้อมร่าง"],
        ],
    )
    add_heading(doc, "7.3 API กลุ่มหลัก (คำนำหน้า /api/v1)", 2)
    add_table(
        doc,
        ["กลุ่ม", "เมธอดตัวอย่าง", "หน้าที่"],
        [
            ["Auth", "POST /auth/login  GET /auth/me", "คุกกี้ JWT"],
            ["Projects", "GET/POST /projects  PATCH .../phase", "วงจรโครงการ"],
            ["Intake", ".../intake/upload analyze coverage confirm-ready chat", "Phase 0–2"],
            ["Draft", "GET/PUT .../sections  POST .../draft-section draft-chat", "Phase 3"],
            ["Review / Export", "POST .../review submit export", "Phase 4"],
            ["Chat / KB", "/chat/rooms  /knowledge-base/mine", "ถาม-ตอบและคลัง"],
            ["Agent", "POST /agent/sessions", "เส้นขนาน (ยังไม่มีหน้า UI)"],
            ["Admin", "GET/PUT /admin/ai-settings  /admin/users", "ระบบ"],
        ],
    )

    add_heading(doc, "8. Rule Engine", 1)
    add_body(doc, "คะแนนรวม = ผลรวมถ่วงน้ำหนัก ผ่านเมื่อ ≥ 70  เรียกซ้ำด้วยอินพุตเดียวกันได้ผลเดียวกัน (ไม่มีสุ่ม)  หัก error −20, warning −10, suggestion −5")
    add_table(
        doc,
        ["หมวด", "น้ำหนัก", "กฎหลัก", "ข้ามเมื่อไร"],
        [
            ["กฎหมาย", "40%", "ทุนจดทะเบียน = floor(งบ÷4) · ค่าปรับ 0.01–0.20%/วัน ขั้นต่ำ 100 บาท · ห้ามล็อกยี่ห้อโดยไม่มี “หรือเทียบเท่า” · อ้าง พ.ร.บ. 2560", "—"],
            ["ความครบ", "30%", "มี 13 หมวด · ยาวขั้นต่ำ · มี s4.1 และ s4.8", "—"],
            ["ความสอดคล้อง", "20%", "งบ↔ขอบเขต · ไทม์ไลน์↔ส่งมอบ · คุณสมบัติ↔ความซับซ้อน", "—"],
            ["รูปแบบ", "10%", "วันที่ไทย พ.ศ. · ลำดับหมวด · รูปแบบราชการ", "—"],
            ["งวดจ่าย (เสริม)", "ร่วมกฎหมาย", "รวมงวด = 100% · งวดละ 5–50%", "ไม่มีข้อมูลงวด"],
            ["ระยะเวลา (เสริม)", "ร่วมกฎหมาย", "งบ > 100 ล้าน → ≥ 180 วัน · งบ < 10 ล้าน → ≤ 365 วัน", "ไม่มีข้อมูลวัน"],
        ],
    )

    add_heading(doc, "9. Tools", 1)
    add_table(
        doc,
        ["บริการ Docker", "บทบาท", "พอร์ต"],
        [
            ["frontend", "Next.js standalone", "3000"],
            ["backend", "FastAPI + Alembic ตอนสตาร์ท", "4000"],
            ["postgres (pgvector)", "ข้อมูลหลัก + เวกเตอร์ 768-d", "5432"],
            ["redis", "แคช เซสชัน อัตราจำกัด", "6379"],
            ["minio", "อ็อบเจ็กต์ DOCX/PDF", "9000 / 9001"],
            ["mongo", "GridFS ต้นฉบับเอกสาร", "27017"],
            ["neo4j", "กราฟกฎหมาย (ลดเหลือ pgvector ถ้าไม่ขึ้น)", "7474 / 7687"],
            ["qdrant", "เวกเตอร์ทางเลือก --profile qdrant", "6333"],
        ],
    )
    add_body(doc, "แชท (LLM_PROVIDER) และ embeddings (EMBEDDING_PROVIDER) เลือกคนละตัวได้ทุกโหมด เช่น Claude API + EmbeddingGemma ใน LM Studio")
    add_table(
        doc,
        ["ชนิด", "ค่าที่รับ"],
        [
            ["โหมด", "on_prem / cloud / hybrid (ป้าย — ไม่บังคับสลับฝั่งอื่น)"],
            ["แชทในเครื่อง", "lm_studio, ollama, llama_cpp"],
            ["แชทคลาวด์", "claude, openai, gemini, bedrock, azure_foundry, openai_compatible"],
            ["Embeddings", "local (ค่าเริ่ม), openai, gemini, bedrock, azure, compat"],
            ["Vector store", "pgvector (ค่าเริ่ม) หรือ qdrant"],
        ],
    )
    add_table(
        doc,
        ["เป้าหมายการตั้งค่า", "แชท", "ฝังเวกเตอร์", "หมายเหตุ"],
        [
            ["ค่าเริ่มต้นในเครื่อง", "LM Studio (Gemma)", "EmbeddingGemma", "โหลดทั้งสองที่พอร์ต 1234"],
            ["Claude + ฝังเวกเตอร์ในเครื่อง", "Claude", "ในเครื่อง", "Anthropic key · ไม่ต้อง seed ใหม่"],
            ["Claude + OpenAI embeddings", "Claude", "OpenAI", "ต้อง seed_raw_docs หลังเปลี่ยนฝังเวกเตอร์"],
        ],
    )
    add_table(
        doc,
        ["คำสั่ง / แพ็ก", "โหลดเข้าแอปเว็บ?"],
        [
            ["python -m app.seed_db", "ใช่ — ผู้ใช้ทดลอง"],
            ["python -m app.seed_raw_docs", "ใช่ — คลังกฎหมายสดจาก PDF"],
            ["python -m app.seed_kb", "ไม่ — extracts งานวิจัย"],
            ["skills/Draft-TORs-Skills", "ไม่ — ร่างออฟไลน์ใน Claude/ChatGPT/Gemini/Hermes"],
            ["skills/check-TORs-Skills", "ไม่ — ตรวจรายการตรวจสอบออฟไลน์"],
        ],
    )

    add_heading(doc, "10. Unit tests", 1)
    add_body(doc, "รันเมื่อ 21 สิงหาคม 2026 บนเครื่องเดียวกันกับรายงานนี้")
    add_image(doc, EVIDENCE / "19-vitest-output.png", "ภาพที่ 15  ผล Vitest — 39 ไฟล์ 167 เคส ผ่าน")
    add_image(doc, EVIDENCE / "19-pytest-output.png", "ภาพที่ 16  ผล pytest — 1500 ผ่าน + live_llm 14")
    add_image(doc, EVIDENCE / "19-unit-test-usage-map.png", "ภาพที่ 17  แผนที่การใช้งาน ↔ เทสต์ที่ล็อกพฤติกรรม")
    add_table(
        doc,
        ["ชุด", "คำสั่ง", "ผลรอบนี้"],
        [
            ["Frontend unit", "cd app/frontend && npm run test:unit", "39 ไฟล์ / 167 เคส ผ่าน"],
            ["Backend unit", 'python -m pytest -m "not live_llm"', "1500 ผ่าน · 0 ข้าม · 86% cov"],
            ["Backend live LLM", "python -m pytest -m live_llm", "14 ผ่าน (LM Studio :1234) รวม realistic workflow"],
            ["E2E headed", "npm run test:e2e:headed", "20 ผ่าน · พิมพ์ช้า + LM Studio จริง + realistic-flow"],
            ["Guide + reports", "npm run test:e2e:guide / reports", "3 + 3 ผ่าน — ดู 18-TEST_EVIDENCE.md"],
        ],
    )
    add_body(doc, "บัญชีทดลองจาก seed_db: officer@example.go.th / Passw0rd!")

    add_heading(doc, "11. สิ่งที่สภาพแวดล้อมนี้ยังไม่ครบ", 1)
    for line in [
        "ถ้าไม่ได้เปิดเซิร์ฟเวอร์ LLM ในเครื่อง Phase 2 / แชทจะโชว์ข้อผิดพลาด — กรอกมือและส่งออกจากข้อความที่มีอยู่ได้",
        "Neo4j ไม่ขึ้น → GraphRAG ลดเหลือ pgvector แชทยังตอบได้แต่ไม่มีกราฟ",
        "ไม่มี e-GP จริง และไม่มี LangGraph checkpointer",
        "/api/v1/agent และ /api/v1/kb-chat ยังไม่มีหน้า Next.js — ใช้ 5 Phase + /chat",
        "Coverage รอบก่อน (19 ส.ค.): backend ~85% · frontend ~89.8%",
    ]:
        add_body(doc, "• " + line)

    add_heading(doc, "12. เอกสารชุดปัจจุบัน", 1)
    add_table(
        doc,
        ["แฟ้ม", "เนื้อหา"],
        [
            ["19 (ต้นฉบับ Markdown ของไฟล์นี้)", "รายงานการทำงานครบ + ภาพ unit tests"],
            ["13-USER_GUIDELINE.md", "คู่มือผู้ใช้ทีละขั้นพร้อมภาพจอ"],
            ["14-INSTALLATION.md", "ติดตั้ง Docker + LM Studio"],
            ["16 / 17", "สถาปัตยกรรม backend / frontend"],
            ["18-TEST_EVIDENCE.md", "หลักฐาน Playwright ครบเคส"],
        ],
    )
    add_body(doc, "เดโม UX ที่คลิกได้โดยไม่เรียก API: https://chiraleo2000.github.io/tor-generator/")
    doc.save(OUT_DOCX)
    print("wrote", OUT_DOCX.name)


def _ppt_font(run, size: int, bold: bool = False, color: PptRgb | None = None) -> None:
    run.font.name = FONT_NAME
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color or P_BODY
    rpr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        el = rpr.find(ppt_qn(f"a:{tag}"))
        if el is not None:
            rpr.remove(el)
        node = rpr.makeelement(ppt_qn(f"a:{tag}"), {"typeface": FONT_NAME})
        rpr.append(node)


def _add_text_box(slide, left, top, width, height, text, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _ppt_font(run, size, bold=bold, color=color)
    return box


def _add_bar(slide) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = P_ORANGE
    shape.line.fill.background()
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(7.2), PptInches(13.333), PptInches(0.3)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = P_NAVY_DARK
    footer.line.fill.background()
    tf = footer.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "TOR Generator  v0.2.3   ·   รายงานการทำงาน   ·   21 ส.ค. 2026"
    _ppt_font(run, 12, color=P_WHITE)


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptRgb(0xFF, 0xFF, 0xFF)
    bg.line.fill.background()
    _add_bar(slide)
    return slide


def _title(slide, text: str) -> None:
    _add_text_box(slide, PptInches(0.5), PptInches(0.28), PptInches(12.3), PptInches(0.5), text, size=26, bold=True, color=P_NAVY)


def _bullets(slide, items: list[str], top=0.95, size=16) -> None:
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(top), PptInches(12.2), PptInches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = PptPt(8)
        run = p.add_run()
        run.text = "•  " + item
        _ppt_font(run, size, color=P_BODY)


def _fill_cell(cell, text: str, size: int, bold: bool, color: PptRgb, fill: PptRgb) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = True
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        _ppt_font(run, size, bold=bold, color=color)


def _table(slide, headers, rows, left=0.5, top=1.05, width=12.3, height=5.4, font=13):
    cols = len(headers)
    table_shape = slide.shapes.add_table(
        1 + len(rows), cols, PptInches(left), PptInches(top), PptInches(width), PptInches(height)
    )
    table = table_shape.table
    for i, header in enumerate(headers):
        _fill_cell(table.cell(0, i), header, font, True, P_WHITE, P_NAVY)
    for r, row in enumerate(rows):
        fill = P_ROW if r % 2 == 0 else P_WHITE
        for c, value in enumerate(row):
            _fill_cell(table.cell(r + 1, c), value, font - 1, False, P_BODY, fill)
    return table


def _picture(slide, path: Path, left, top, width) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), PptInches(left), PptInches(top), width=PptInches(width))


def build_pptx(arch: Path, phases: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = P_NAVY_DARK
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.18), PptInches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = P_ORANGE
    accent.line.fill.background()
    _add_text_box(s, PptInches(0.7), PptInches(1.8), PptInches(11.8), PptInches(0.5), "รายงานการทำงานของแอปพลิเคชัน", size=22, color=P_ORANGE)
    _add_text_box(s, PptInches(0.7), PptInches(2.3), PptInches(11.8), PptInches(0.8), "TOR Generator", size=44, bold=True, color=P_WHITE)
    _add_text_box(
        s,
        PptInches(0.7),
        PptInches(3.2),
        PptInches(11.8),
        PptInches(0.8),
        "ระบบจัดทำและตรวจสอบ TOR ตาม พ.ร.บ. การจัดซื้อจัดจ้างและการบริหารพัสดุภาครัฐ พ.ศ. 2560",
        size=18,
        color=P_WHITE,
    )
    _add_text_box(
        s,
        PptInches(0.7),
        PptInches(5.6),
        PptInches(11.8),
        PptInches(0.6),
        "เวอร์ชัน 0.2.3   ·   21 สิงหาคม 2026   ·   Next.js 14 + FastAPI + LangGraph",
        size=16,
        color=P_ORANGE,
    )

    s = _blank(prs)
    _title(s, "สารบัญ")
    _table(
        s,
        ["บท", "เนื้อหา"],
        [
            ["1–2", "สรุปหนึ่งหน้า · ผู้ใช้และสิทธิ์"],
            ["3–4", "Frontend Next.js · Backend FastAPI + 13 หมวด TOR"],
            ["5", "Workflows 5 Phase, LangGraph, แชท, ตรวจไฟล์, อนุมัติ"],
            ["6–7", "Features และ Functions ที่ขับการใช้งาน"],
            ["8–9", "Rule Engine · Docker / LLM / seed / skills"],
            ["10", "Unit tests จริง 21 ส.ค. 2026 (Vitest 167 / pytest 1500 + live 14 / E2E headed 20)"],
        ],
        height=4.8,
    )

    s = _blank(prs)
    _title(s, "แอปนี้แก้ปัญหาอะไร")
    _bullets(
        s,
        [
            "เขียน TOR ด้วยมือใช้เวลานาน ตกหัวข้อกฎหมาย และจัดรูปแบบราชการยาก",
            "บังคับโครง 13 ส่วน (s1–s13) + ขอบเขตงานย่อย s4.1–s4.14 รวม 27 ช่อง",
            "รับชุดเอกสารในแชท → จัดเข้าช่อง → ตรวจด้วย Rule Engine → คนยืนยันหมวดเสี่ยง → ส่งออก Word/PDF",
            "เส้นทางหลักบนหน้าจอคือ 5 Phase ไม่ใช่วิซาร์ด 8 ขั้น (เหลือไว้เพื่อความเข้ากันได้)",
            "ค่าเริ่มต้นประมวลผลในเครื่อง: Gemma + EmbeddingGemma ผ่าน LM Studio ไม่ส่งข้อมูลออกถ้าอยู่โหมด on_prem",
        ],
    )

    s = _blank(prs)
    _title(s, "สถาปัตยกรรมรวม")
    _picture(s, arch, 0.45, 0.95, 12.4)

    s = _blank(prs)
    _title(s, "ผู้ใช้และสิทธิ์")
    _table(
        s,
        ["บทบาท", "ทำได้", "ทำไม่ได้"],
        [
            ["เจ้าหน้าที่", "โครงการของตนเอง, คลังส่วนตัว, ถาม-ตอบ, ส่งตรวจ", "เห็นของคนอื่น / อนุมัติ / ตั้งค่า AI"],
            ["ผู้ตรวจสอบ", "เห็นทุกโครงการ, อนุมัติหรือส่งกลับ", "แม่แบบ ผู้ใช้ การตั้งค่า AI"],
            ["ผู้ดูแล", "ทุกอย่างของ reviewer + แอดมินคลัสเตอร์", "—"],
        ],
        height=2.4,
        font=14,
    )
    _add_text_box(
        s,
        PptInches(0.55),
        PptInches(3.7),
        PptInches(12.2),
        PptInches(2.8),
        "JWT อยู่ในคุกกี้ HttpOnly tor_access_token (SameSite=Lax)\n"
        "เมนูหลัก: แดชบอร์ด · ฐานความรู้ · ร่าง TOR · ตรวจสอบ TOR · ถาม-ตอบ · คู่มือ\n"
        "แอดมิน: แม่แบบ · ฐานความรู้ (จัดการ) · ผู้ใช้ · การตั้งค่า AI",
        size=16,
    )

    s = _blank(prs)
    _title(s, "หน้าจอหลัก")
    _picture(s, EVIDENCE / "00-login.png", 0.4, 1.0, 6.1)
    _picture(s, EVIDENCE / "02-dashboard.png", 6.7, 1.0, 6.1)
    _add_text_box(s, PptInches(0.4), PptInches(6.55), PptInches(6.1), PptInches(0.4), "เข้าสู่ระบบ", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)
    _add_text_box(s, PptInches(6.7), PptInches(6.55), PptInches(6.1), PptInches(0.4), "แดชบอร์ดโครงการ", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)

    s = _blank(prs)
    _title(s, "สามเครื่องมือหลักของเจ้าหน้าที่")
    _table(
        s,
        ["เครื่องมือ", "หน้า", "ผลที่ควรเห็น"],
        [
            ["ร่าง TOR", "/projects/{id}/draft · 5 Phase", "เนื้อหา s1–s13 พร้อมส่งตรวจ"],
            ["ตรวจสอบ TOR", "Phase 4 + /review", "คะแนน ≥70 / findings / suggestions"],
            ["ถาม-ตอบ", "/chat (SSE + citations)", "คำตอบยาวพร้อมชิปอ้างอิงจากคลัง"],
        ],
        height=3.0,
        font=16,
    )
    _picture(s, EVIDENCE / "13-kb-chat.png", 0.55, 4.3, 5.8)
    _picture(s, EVIDENCE / "12-standalone-review.png", 6.6, 4.3, 5.8)

    s = _blank(prs)
    _title(s, "Frontend — Next.js 14")
    _table(
        s,
        ["ชั้น", "รายละเอียด"],
        [
            ["คำขอ", "Browser → /api/v1 (โดเมนเดียว) → rewrite backend:4000 timeout 5 นาที"],
            ["สถานะ", "Zustand: auth / project / ui  ·  JWT ไม่อยู่ใน localStorage"],
            ["ร่าง", "DraftWorkspace คือ UI จริง  ·  วิซาร์ด 8 ขั้นเป็นเส้นเข้ากันได้"],
            ["เกตเฟส", "phase-gate.ts วิเคราะห์แล้วเลือก Phase 2 ได้ · ready_to_compose ปลด Phase 3"],
            ["แชท", "ChatShell + SSE  ·  kind=kb ที่ /chat  คนละประวัติกับ draft_intake"],
        ],
        height=4.6,
        font=14,
    )

    s = _blank(prs)
    _title(s, "Backend — ชั้นระบบ")
    _table(
        s,
        ["ชั้น", "โฟลเดอร์", "หน้าที่"],
        [
            ["HTTP", "api/v1/endpoints", "REST + SSE"],
            ["Domain", "domain", "s1–s13, slots, magic bytes"],
            ["Orchestrator", "orchestrator", "กราฟรายหมวด + กราฟเอเจนต์"],
            ["RAG", "rag", "สกัด หั่น ค้น ACL กราฟ"],
            ["Rules", "rule_engine", "คะแนนกำหนดได้ซ้ำ ผ่านที่ ≥ 70"],
            ["Providers", "providers", "LLM / embeddings / vector อิสระ"],
            ["Export", "export", "DOCX PDF ไทย MinIO"],
        ],
        height=5.2,
        font=13,
    )

    s = _blank(prs)
    _title(s, "13 หมวด TOR และจุด HITL")
    _table(
        s,
        ["คีย์", "ชื่อ", "ต้องคนยืนยัน"],
        [
            ["s1 / s2", "ความเป็นมา / วัตถุประสงค์", "ไม่"],
            ["s3", "คุณสมบัติผู้เสนอราคา", "ใช่"],
            ["s4", "ขอบเขตงาน + s4.1–s4.14", "ไม่ (บังคับย่อย s4.1, s4.8)"],
            ["s5 / s7", "ระยะเวลา / สถานที่", "ไม่"],
            ["s6 / s8", "งบประมาณ / งวดจ่าย", "ใช่"],
            ["s9 / s11 / s12", "รับประกัน / เกณฑ์ / เอกสารยื่น", "ไม่"],
            ["s10 / s13", "ค่าปรับ / เงื่อนไขอื่น", "ใช่"],
        ],
        height=5.2,
        font=14,
    )

    s = _blank(prs)
    _title(s, "Workflow หลัก — 5 Phase")
    _picture(s, phases, 0.4, 1.0, 12.5)

    s = _blank(prs)
    _title(s, "เกตแต่ละเฟส")
    _table(
        s,
        ["Phase", "งานผู้ใช้", "ปลดล็อกเมื่อ"],
        [
            ["0 อัปโหลด", "ลากไฟล์เข้าแชทโครงการ", "มีไฟล์หรือข้อความ"],
            ["1 ช่องว่าง", "ตารางความครบ ดึงกฎหมาย กดพร้อมร่าง", "ready_to_compose = true"],
            ["2 ร่าง", "แก้หรือ AI ร่าง 13 หมวด + HITL", "ครบ 13 และยืนยัน s3 s6 s8 s10 s13"],
            ["3 ตรวจ", "คะแนน Rule Engine แล้วส่งขออนุมัติ", "reviewer/admin อนุมัติหรือส่งกลับ"],
            ["4 ส่งออก", "Word / PDF ฟอนต์ไทย วันที่ พ.ศ.", "e-Bidding อยู่นอกแอป"],
        ],
        height=5.2,
        font=14,
    )

    s = _blank(prs)
    _title(s, "ภาพการใช้งาน Phase 0–2")
    _picture(s, EVIDENCE / "03-phase-0-upload.png", 0.35, 0.95, 6.2)
    _picture(s, EVIDENCE / "04b-phase-1-coverage.png", 6.8, 0.95, 6.1)
    _add_text_box(s, PptInches(0.35), PptInches(6.5), PptInches(6.2), PptInches(0.4), "Phase 0 อัปโหลด", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)
    _add_text_box(s, PptInches(6.8), PptInches(6.5), PptInches(6.1), PptInches(0.4), "Phase 1 ความครบ", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)

    s = _blank(prs)
    _title(s, "ภาพการใช้งาน Phase 2–4")
    _picture(s, EVIDENCE / "05-phase-2-draft.png", 0.3, 0.95, 4.2)
    _picture(s, EVIDENCE / "06-phase-3-review.png", 4.55, 0.95, 4.2)
    _picture(s, EVIDENCE / "07-phase-4-publish.png", 8.8, 0.95, 4.2)
    _add_text_box(s, PptInches(0.3), PptInches(6.5), PptInches(4.2), PptInches(0.4), "ร่าง + HITL", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)
    _add_text_box(s, PptInches(4.55), PptInches(6.5), PptInches(4.2), PptInches(0.4), "ตรวจคะแนน", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)
    _add_text_box(s, PptInches(8.8), PptInches(6.5), PptInches(4.2), PptInches(0.4), "ส่งออก", size=13, color=P_MUTED, align=PP_ALIGN.CENTER)

    s = _blank(prs)
    _title(s, "LangGraph และเส้นขนานเอเจนต์")
    _bullets(
        s,
        [
            "ร่างรายหมวด: validate → RAG → llm_draft → rule_guardrail (≥70 หรือ retry ≤3) → HITL → finalize",
            "Timeout ต่อหมวด 180 วินาทีสำหรับ Gemma (สูงสุด 300)",
            "เส้น /api/v1/agent: ingest → map 27 ช่อง → ถามส่วนขาด → ยืนยัน → ร่างทั้งฉบับ → HITL → ส่งออก",
            "เส้นเอเจนต์ยังไม่มีหน้า Next.js — UI จริงคือ 5 Phase + /chat",
            "ถาม-ตอบคลังที่ /chat (kind=kb) คนละประวัติกับแชทร่าง (kind=draft_intake)",
            "ตรวจไฟล์ภายนอกที่ /review ด้วย Jaccard โดยไม่ต้องสร้างโครงการ",
        ],
        size=16,
    )

    s = _blank(prs)
    _title(s, "Features หลัก")
    _table(
        s,
        ["ฟีเจอร์", "จุดสำคัญ"],
        [
            ["Intake หลายไฟล์", "ไม่ต้องเลือก 9 ประเภท · magic bytes · GridFS"],
            ["แผนที่ช่อง", "NLP → s1–s13 / s4.1–s4.14"],
            ["ร่างด้วย AI", "เอเจนต์ 13 ตัว + RAG + Rule Engine"],
            ["HITL", "s3 คุณสมบัติ, s6 งบ, s8 งวด, s10 ค่าปรับ, s13 เงื่อนไข"],
            ["คลังความรู้", "คลังกลาง + เอกสารของฉัน (ACL ตาม owner_id)"],
            ["ตั้งค่า AI", "แชทกับ embeddings เลือกอิสระ มีผลทันทีหลังบันทึก"],
        ],
        height=5.2,
        font=14,
    )

    s = _blank(prs)
    _title(s, "Rule Engine — น้ำหนักคะแนน ผ่านที่ ≥ 70")
    _table(
        s,
        ["หมวด", "น้ำหนัก", "ตัวอย่างกฎ"],
        [
            ["กฎหมาย", "40%", "ทุนจดทะเบียน floor(งบ÷4) · ค่าปรับ 0.01–0.20%/วัน"],
            ["ความครบ", "30%", "ครบ 13 หมวด · มี s4.1 และ s4.8"],
            ["ความสอดคล้อง", "20%", "งบ↔ขอบเขต · ไทม์ไลน์↔ส่งมอบ"],
            ["รูปแบบ", "10%", "วันที่ พ.ศ. · ลำดับหมวดราชการ"],
            ["งวดจ่าย / ระยะเวลา", "ร่วมกฎหมาย", "ข้ามเมื่อไม่มีข้อมูล"],
        ],
        height=4.4,
        font=14,
    )
    _add_text_box(s, PptInches(0.55), PptInches(5.7), PptInches(12), PptInches(0.8), "หักคะแนน: error −20  ·  warning −10  ·  suggestion −5   ·   ผลซ้ำได้ 100% (ไม่มีสุ่ม)", size=16, color=P_NAVY)

    s = _blank(prs)
    _title(s, "Tools — Docker, LLM, seed, skills")
    _table(
        s,
        ["เครื่องมือ", "บทบาท"],
        [
            ["Docker", "frontend 3000, backend 4000, Postgres, Redis, MinIO, Mongo, Neo4j"],
            ["LLM ค่าเริ่ม", "lm_studio + google/gemma-4-e4b · EmbeddingGemma 768-d"],
            ["คลาวด์ที่รองรับ", "Claude, OpenAI, Gemini, Bedrock, Azure Foundry, OpenAI-compat"],
            ["seed_raw_docs", "คลังกฎหมายสดจาก PDF ใน documents/sources/"],
            ["skills/", "แพ็กออฟไลน์ Draft-TORs / check-TORs — ไม่ถูกโหลดโดยเว็บ"],
        ],
        height=5.2,
        font=14,
    )

    s = _blank(prs)
    _title(s, "Functions ที่ล็อกพฤติกรรมผู้ใช้")
    _table(
        s,
        ["ฟังก์ชัน", "ผลต่อการใช้งาน"],
        [
            ["canSelectPhase", "วิเคราะห์แล้วเลือก Phase 2 ได้ · Phase 3 ต้อง ready_to_compose"],
            ["validatePassword", "สมัครได้เมื่อรหัสผ่านครบกฎ (ตัวเลขต้องเป็น ASCII)"],
            ["streamSsePost", "แชทคลังและแชทร่างสตรีมคำตอบ"],
            ["get_agent_for_section", "ปุ่มร่างด้วย AI เรียกเอเจนต์ถูกหมวด"],
            ["RuleEngine.validate", "Phase 4 ให้คะแนนตาม พ.ร.บ. 2560"],
            ["ProviderFactory", "แอดมินผสมแชทคลาวด์กับ embeddings ในเครื่องได้"],
            ["require_project_access", "เจ้าหน้าที่ไม่เห็นโครงการหรือคลังของคนอื่น"],
        ],
        height=5.3,
        font=13,
    )

    s = _blank(prs)
    _title(s, "Unit tests ที่รันจริง — 21 ส.ค. 2026")
    _picture(s, EVIDENCE / "19-vitest-output.png", 0.35, 0.95, 6.2)
    _picture(s, EVIDENCE / "19-pytest-output.png", 6.75, 0.95, 6.2)

    s = _blank(prs)
    _title(s, "เทสต์ล็อกการใช้งานอย่างไร")
    _picture(s, EVIDENCE / "19-unit-test-usage-map.png", 0.55, 0.9, 12.2)

    s = _blank(prs)
    _title(s, "วิธีรันเทสต์เอง")
    _bullets(
        s,
        [
            "Frontend:  cd app/frontend && npm run test:unit",
            'Backend:   cd app/backend && python -m pytest -m "not live_llm" --hypothesis-profile=coverage',
            "ถ้า LM Studio เปิดที่ :1234:  python -m pytest -m live_llm",
            "E2E headed (ต้องมี UI ที่ :3000):  npm run test:e2e:headed",
            "Guide screenshots:  npm run test:e2e:guide",
            "บัญชีทดลอง: officer@example.go.th / Passw0rd!",
            "live_llm 14 เคสผ่านเมื่อเปิด LM Studio ที่ :1234 — รวม backend 1514 เคส",
        ],
        size=17,
    )

    s = _blank(prs)
    _title(s, "ขอบเขตที่ยังไม่ครบในสภาพแวดล้อมนี้")
    _bullets(
        s,
        [
            "ไม่มี LLM ในเครื่อง → Phase 2 / แชทผิดพลาด แต่กรอกมือและส่งออกได้",
            "Neo4j ไม่ขึ้น → แชทยังตอบจาก pgvector แต่ไม่มีกราฟกฎหมาย",
            "ไม่มี e-GP จริง และไม่มี LangGraph checkpointer",
            "/api/v1/agent และ /api/v1/kb-chat ยังไม่มีหน้า Next.js",
            "อย่าใช้เอกสารออกแบบ 10 / 11 เป็นคู่มือติดตั้งของสแตก Docker ปัจจุบัน",
        ],
    )

    s = _blank(prs)
    _title(s, "สรุป")
    _bullets(
        s,
        [
            "เจ้าหน้าที่เดิน 5 Phase: อัปโหลด → เติมช่อง → ร่าง 13 หมวด → ตรวจ → ส่งออก",
            "คนยืนยันหมวดเสี่ยง (HITL) ก่อนส่งขออนุมัติ ผู้ตรวจสอบกดอนุมัติหรือส่งกลับ",
            "AI ในเครื่องเป็นค่าเริ่ม แชทกับ embeddings เลือกอิสระได้จากหน้าแอดมิน",
            "Unit tests ล็อกพฤติกรรมนี้ไว้แล้ว: Vitest 167, pytest 1500 + live_llm 14, Playwright headed 20",
            "ต้นฉบับ Markdown: Discussions/19-APPLICATION_OPERATING_REPORT.md",
        ],
        size=18,
    )

    prs.save(OUT_PPTX)
    print("wrote", OUT_PPTX.name)


def docx_to_pdf_word() -> bool:
    script = f"""
$ErrorActionPreference = 'Stop'
$word = New-Object -ComObject Word.Application
$word.Visible = $false
$word.DisplayAlerts = 0
try {{
  $docx = '{str(OUT_DOCX).replace("'", "''")}'
  $pdf = '{str(OUT_PDF).replace("'", "''")}'
  $doc = $word.Documents.Open($docx, $false, $true)
  if (Test-Path $pdf) {{ Remove-Item $pdf -Force }}
  $wdExportFormatPDF = 17
  $doc.ExportAsFixedFormat($pdf, $wdExportFormatPDF)
  $doc.Close($false)
  Write-Output 'PDF_OK'
}} finally {{
  $word.Quit()
  [System.Runtime.InteropServices.Marshal]::ReleaseComObject($word) | Out-Null
}}
"""
    result = subprocess.run(
        ["powershell", "-NoProfile", "-Command", script],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    sys.stdout.write(result.stdout)
    sys.stderr.write(result.stderr)
    return result.returncode == 0 and OUT_PDF.exists()


def docx_to_pdf_edge() -> bool:
    html = ROOT / "_report19_print.html"
    html.write_text(
        f"""<!doctype html><meta charset="utf-8"><title>Redirect</title>
<p>Use Word export. This fallback prints a notice.</p>""",
        encoding="utf-8",
    )
    html.unlink(missing_ok=True)
    return False


def main() -> None:
    EVIDENCE.mkdir(exist_ok=True)
    arch = draw_architecture()
    phases = draw_phases()
    print("diagrams", arch.name, phases.name)
    build_docx(arch, phases)
    build_pptx(arch, phases)
    if docx_to_pdf_word():
        print("wrote", OUT_PDF.name)
        return
    print("Word COM PDF failed; trying nothing further would lose Thai layout")
    raise SystemExit(1)


if __name__ == "__main__":
    main()
